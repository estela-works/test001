"""各スライドPNG画像をPowerPointに全面配置して保存する"""
import os
from pptx import Presentation
from pptx.util import Emu

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SLIDE_DIR = os.path.dirname(SCRIPT_DIR)
IMAGES_DIR = os.path.join(SLIDE_DIR, 'images')
OUTPUT_PATH = os.path.join(SLIDE_DIR, 'presentation.pptx')

# スライド順序（slide-viewer.html のインラインJSONと同じ順序）
SLIDE_NAMES = [
    'main-0-01', 'main-0-02', 'main-0-03', 'main-0-04', 'main-0-05', 'main-0-06',
    'main-1-01', 'main-1-02', 'main-1-03', 'main-1-04', 'main-1-05',
    'main-1-06', 'main-1-07', 'main-1-08', 'main-1-09', 'main-1-10',
    'main-2-01', 'main-2-02', 'main-2-03', 'main-2-04', 'main-2-05',
    'main-2-06', 'main-2-07', 'main-2-08', 'main-2-09',
    'main-3-01', 'main-3-02', 'main-3-04', 'main-3-05', 'main-3-06',
    'main-4-01', 'main-4-02', 'main-4-03', 'main-4-04', 'main-4-07', 'main-4-08',
    'main-5-01', 'main-5-02', 'main-5-03', 'main-5-04', 'main-5-05',
    'main-5-07', 'main-5-09', 'main-5-10',
    'main-6-01', 'main-6-03', 'main-6-04',
    'main-c-01', 'main-c-02', 'main-c-03', 'main-c-04',
]

# 16:9 ワイドスライド (25.4cm x 14.29cm)
SLIDE_WIDTH = Emu(9144000)   # 25.4cm
SLIDE_HEIGHT = Emu(5143500)  # 14.29cm

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # blank layout

for name in SLIDE_NAMES:
    img_path = os.path.join(IMAGES_DIR, f'{name}.png')
    if not os.path.exists(img_path):
        print(f'SKIP (not found): {name}')
        continue

    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    print(f'added: {name}')

prs.save(OUTPUT_PATH)
print(f'\nSaved: {OUTPUT_PATH} ({len(prs.slides)} slides)')
